from langchain_community.document_loaders import PyMuPDFLoader, TextLoader
from pptx import Presentation
import fitz
from langchain_classic.schema import Document
import os
from PIL import Image
from io import BytesIO
from docx import Document as DocxDocument

def is_scanned(filepath):
    doc = fitz.open(filepath)
    for page in doc:
        if page.get_text().strip():
            return False
    return True


def scanned_pdf_loader(filepath):
    pdf = fitz.open(filepath)
    docs = []

    for i, page in enumerate(pdf):
        pix = page.get_pixmap(dpi=150)
        img = Image.open(BytesIO(pix.tobytes("png")))

        text = page.get_text("text") 
        docs.append(
            Document(
                page_content=text,
                metadata={"source": filepath, "page": i, "ocr": True},
            )
        )

    return docs


def load_pptx(filepath):
    ppt = Presentation(filepath)
    docs = []

    for i, slide in enumerate(ppt.slides):
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]).strip()
        if text:
            docs.append(Document(page_content=text, metadata={"source": filepath, "page": i}))
    return docs


def load_docx(filepath):
    doc = DocxDocument(filepath)
    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    return [Document(page_content=text, metadata={"source": filepath})]


def load_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pdf":
        return scanned_pdf_loader(filepath) if is_scanned(filepath) else PyMuPDFLoader(filepath).load()

    if ext == ".docx":
        return load_docx(filepath)

    if ext == ".pptx":
        return load_pptx(filepath)

    if ext == ".txt":
        return TextLoader(filepath).load()

    raise ValueError(f"Unsupported file format: {ext}")
